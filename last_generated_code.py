from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# 创建演示文稿
prs = Presentation()

# 第一页：封面页
slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # 标题幻灯片布局
title = slide1.shapes.title
subtitle = slide1.placeholders[1]

title.text = "周八公司介绍汇报"
subtitle.text = "企业概况与发展规划"
title.text_frame.paragraphs[0].font.size = Pt(44)
title.text_frame.paragraphs[0].font.bold = True
subtitle.text_frame.paragraphs[0].font.size = Pt(28)

# 第二页：公司概况
slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # 标题和内容布局
title2 = slide2.shapes.title
content2 = slide2.placeholders[1]

title2.text = "公司概况"
title2.text_frame.paragraphs[0].font.size = Pt(36)
title2.text_frame.paragraphs[0].font.bold = True

tf = content2.text_frame
tf.clear()  # 清除默认文本

p = tf.add_paragraph()
p.text = "企业基本信息"
p.font.size = Pt(24)
p.font.bold = True

p = tf.add_paragraph()
p.text = "• 公司名称：周八科技有限公司"
p.font.size = Pt(20)
p.level = 0

p = tf.add_paragraph()
p.text = "• 成立时间：2018年"
p.font.size = Pt(20)
p.level = 0

p = tf.add_paragraph()
p.text = "• 注册资本：5000万元"
p.font.size = Pt(20)
p.level = 0

p = tf.add_paragraph()
p.text = "• 员工规模：200+人"
p.font.size = Pt(20)
p.level = 0

p = tf.add_paragraph()
p.text = "• 主营业务：软件开发、技术服务、数字化转型解决方案"
p.font.size = Pt(20)
p.level = 0

p = tf.add_paragraph()
p.text = ""
p.font.size = Pt(20)

p = tf.add_paragraph()
p.text = "企业愿景"
p.font.size = Pt(24)
p.font.bold = True

p = tf.add_paragraph()
p.text = "致力于成为国内领先的数字化转型服务提供商，为客户创造更大价值"
p.font.size = Pt(20)
p.level = 0

# 第三页：发展规划
slide3 = prs.slides.add_slide(prs.slide_layouts[1])  # 标题和内容布局
title3 = slide3.shapes.title
content3 = slide3.placeholders[1]

title3.text = "发展规划"
title3.text_frame.paragraphs[0].font.size = Pt(36)
title3.text_frame.paragraphs[0].font.bold = True

tf3 = content3.text_frame
tf3.clear()

p = tf3.add_paragraph()
p.text = "短期目标（2024-2025）"
p.font.size = Pt(24)
p.font.bold = True

p = tf3.add_paragraph()
p.text = "• 完成新一轮融资，扩大市场份额"
p.font.size = Pt(20)
p.level = 0

p = tf3.add_paragraph()
p.text = "• 推出3款核心产品，提升技术壁垒"
p.font.size = Pt(20)
p.level = 0

p = tf3.add_paragraph()
p.text = "• 建立完善的服务体系，提升客户满意度"
p.font.size = Pt(20)
p.level = 0

p = tf3.add_paragraph()
p.text = ""
p.font.size = Pt(20)

p = tf3.add_paragraph()
p.text = "长期愿景（2026-2030）"
p.font.size = Pt(24)
p.font.bold = True

p = tf3.add_paragraph()
p.text = "• 成为行业独角兽企业，估值超10亿元"
p.font.size = Pt(20)
p.level = 0

p = tf3.add_paragraph()
p.text = "• 拓展海外市场，实现国际化布局"
p.font.size = Pt(20)
p.level = 0

p = tf3.add_paragraph()
p.text = "• 构建产业生态圈，引领行业发展方向"
p.font.size = Pt(20)
p.level = 0

# 保存PPT到当前目录
ppt_path = "周八公司介绍汇报.pptx"
prs.save(ppt_path)

# 获取当前工作目录
current_dir = os.getcwd()
full_path = os.path.join(current_dir, ppt_path)

print(f"PPT已生成：{full_path}")
print(f"文件大小：{os.path.getsize(ppt_path)} 字节")