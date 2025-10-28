"""
正确使用模板的示例代码
这个代码会保留模板的背景、logo和所有设计元素
"""
from pptx import Presentation

# ✅ 正确：加载模板（保留背景、logo、配色）
template_path = 'skills_storage/pptx/temple/附件2：进门_转正答辩PPT模板.pptx'
prs = Presentation(template_path)

print(f"模板总共有 {len(prs.slides)} 张幻灯片")

# 确定需要的页数（不超过模板现有页数）
needed_slides = min(4, len(prs.slides))  # 需要4页，但不超过模板现有数量

# ✅ 正确：删除多余的幻灯片（从后往前删除）
total_slides = len(prs.slides)
if total_slides > needed_slides:
    for i in range(total_slides - 1, needed_slides - 1, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

print(f"删除后剩余 {len(prs.slides)} 张幻灯片")

# ✅ 正确：修改现有幻灯片的文本内容
# 第1页 - 封面
slide_0 = prs.slides[0]
for shape in slide_0.shapes:
    if hasattr(shape, "text_frame"):
        # 识别标题位置（通常在上方）
        if shape.top < 2000000:  # 上方区域
            if len(shape.text) > 0 or shape.text_frame.text.strip() == "":
                shape.text = "王五公司汇报"
        elif 2000000 <= shape.top < 4000000:  # 中间区域
            shape.text = "2024年度工作总结"

# 第2页 - 公司概况
slide_1 = prs.slides[1]
title_set = False
for shape in slide_1.shapes:
    if hasattr(shape, "text_frame"):
        if shape.top < 1500000 and not title_set:  # 标题区域
            shape.text = "公司概况"
            title_set = True
        elif shape.top >= 1500000:  # 内容区域
            tf = shape.text_frame
            tf.clear()
            tf.text = "• 公司名称：王五有限公司"
            p = tf.add_paragraph()
            p.text = "• 成立时间：2010年"
            p.level = 0
            p = tf.add_paragraph()
            p.text = "• 员工人数：500+"
            p.level = 0
            p = tf.add_paragraph()
            p.text = "• 主营业务：技术服务"
            p.level = 0

# 第3页 - 业务发展与未来展望
if len(prs.slides) > 2:
    slide_2 = prs.slides[2]
    title_set = False
    for shape in slide_2.shapes:
        if hasattr(shape, "text_frame"):
            if shape.top < 1500000 and not title_set:
                shape.text = "业务发展与未来展望"
                title_set = True
            elif shape.top >= 1500000:
                tf = shape.text_frame
                tf.clear()
                tf.text = "业务发展："
                p = tf.add_paragraph()
                p.text = "• 年度业绩：营收增长35%"
                p.level = 0
                p = tf.add_paragraph()
                p.text = "• 客户数量：超过1000家"
                p.level = 0
                p = tf.add_paragraph()
                p.text = ""
                p = tf.add_paragraph()
                p.text = "未来展望："
                p = tf.add_paragraph()
                p.text = "• 战略目标：行业领先"
                p.level = 0
                p = tf.add_paragraph()
                p.text = "• 技术投入：AI、大数据"
                p.level = 0

# 保存（会保留所有模板设计）
output_file = '王五公司汇报_正确版本.pptx'
prs.save(output_file)
print(f"\n✅ PPT 已生成：{output_file}")
print("✅ 背景、logo、配色已保留")

