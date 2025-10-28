"""
测试脚本：验证模板背景和logo是否能正确保留
"""
from pptx import Presentation

# 加载模板
template_path = 'skills_storage/pptx/temple/附件2：进门_转正答辩PPT模板.pptx'
prs = Presentation(template_path)

print(f"模板总共有 {len(prs.slides)} 张幻灯片")

# 只保留前4页，删除其他
total_slides = len(prs.slides)
for i in range(total_slides - 1, 3, -1):  # 从后往前删除第5页及之后的
    rId = prs.slides._sldIdLst[i].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[i]

print(f"删除后剩余 {len(prs.slides)} 张幻灯片")

# 修改文本内容（保留所有背景、logo等设计元素）
for slide_idx, slide in enumerate(prs.slides):
    print(f"\n--- 幻灯片 {slide_idx} ---")
    for shape_idx, shape in enumerate(slide.shapes):
        if hasattr(shape, "text_frame"):
            print(f"  形状 {shape_idx}: 位置(top={shape.top}, left={shape.left}), 文本='{shape.text[:30]}'")
            
            # 简单修改：在第一页修改标题
            if slide_idx == 0 and shape_idx == 0:
                shape.text = "测试：保留背景和Logo"
            elif slide_idx == 0 and shape_idx == 1:
                shape.text = "验证模板是否正确使用"

# 保存
output_file = 'test_template_output.pptx'
prs.save(output_file)
print(f"\n✅ 测试完成！请打开 {output_file} 检查：")
print("   1. 背景是否保留？")
print("   2. Logo是否保留？")
print("   3. 配色方案是否保留？")

